import os
import re
import asyncio
from pathlib import Path
from typing import Optional, List, Tuple
from datetime import datetime
import aiofiles

from backend.config import BRAIN_INBOX, BRAIN_NOTES


IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif"}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".html", ".htm", ".txt", ".csv",
                       ".json", ".xml", ".rtf", ".odt", ".epub"}

_OPTIMIZED_DIR = "img"


def optimize_image_sync(src_path: str, max_dim: int = 1920, quality: int = 85,
                        subfolder: str = "") -> Optional[Tuple[str, dict]]:
    """Optimize an image: resize, convert to WebP, strip EXIF.

    Saves to brain/baul/img/[subfolder/]<stem>.webp
    Returns (relative_path, metadata_dict) or None on failure.
    """
    try:
        from PIL import Image

        img = Image.open(src_path)
        original = {"width": img.width, "height": img.height, "format": img.format, "mode": img.mode, "size": os.path.getsize(src_path)}

        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")

        if img.width > max_dim or img.height > max_dim:
            ratio = min(max_dim / img.width, max_dim / img.height)
            new_w = int(img.width * ratio)
            new_h = int(img.height * ratio)
            img = img.resize((new_w, new_h), Image.LANCZOS)
            original["resized"] = True
            original["original_size"] = (original["width"], original["height"])
            original["width"] = new_w
            original["height"] = new_h
        else:
            original["resized"] = False

        stem = Path(src_path).stem
        if subfolder:
            subfolder_clean = subfolder.strip("/\\").replace("\\", "/")
            out_dir = BRAIN_NOTES / _OPTIMIZED_DIR / subfolder_clean
        else:
            out_dir = BRAIN_NOTES / _OPTIMIZED_DIR
        out_dir.mkdir(parents=True, exist_ok=True)
        out_path = out_dir / f"{stem}.webp"

        img.save(str(out_path), "WEBP", quality=quality)
        original["optimized_size"] = out_path.stat().st_size
        original["optimized_path"] = str(out_path.relative_to(BRAIN_NOTES)).replace("\\", "/")
        original["format"] = "webp"
        original["savings_pct"] = round((1 - original["optimized_size"] / max(original["size"], 1)) * 100, 1)

        return str(out_path.relative_to(BRAIN_NOTES)).replace("\\", "/"), original
    except Exception as e:
        return None


def _ocr_image_sync(image_path: str) -> str:
    """Extract text from image using EasyOCR."""
    import sys
    import io
    # Preserve and restore stdout/stderr to handle encoding issues
    old_out = sys.stdout
    old_err = sys.stderr
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')
    try:
        import easyocr
        reader = easyocr.Reader(["en", "es"], gpu=False, verbose=False)
        results = reader.readtext(image_path, detail=0, paragraph=True)
        if results:
            cleaned = []
            for r in results:
                clean = r.strip()
                if clean:
                    cleaned.append(clean)
            return "\n\n".join(cleaned)
        return ""
    except Exception as e:
        raise Exception(f"OCR failed: {e}")
    finally:
        sys.stdout = old_out
        sys.stderr = old_err


def _image_metadata_sync(image_path: str) -> dict:
    """Extract image metadata using Pillow."""
    from PIL import Image
    img = Image.open(image_path)
    meta = {
        "width": img.width,
        "height": img.height,
        "format": img.format,
        "mode": img.mode,
    }
    img.close()
    return meta


def _is_likely_image_bytes(filepath: Path) -> bool:
    """Detect if file is an image by reading magic bytes, regardless of extension."""
    try:
        with open(filepath, "rb") as f:
            header = f.read(12)
        if header[:4] == b"\x89PNG": return True
        if header[:2] == b"\xff\xd8": return True
        if header[:3] in (b"GIF",): return True
        if header[:2] == b"BM": return True
        if header[:4] == b"RIFF" and header[8:12] == b"WEBP": return True
        if header[:4] == b"\x49\x49\x2a\x00": return True
        if header[:4] == b"\x4d\x4d\x00\x2a": return True
        return False
    except Exception:
        return False


def _is_likely_text_document(filepath: Path) -> bool:
    """Check if a file is likely a text document vs an image."""
    ext = filepath.suffix.lower()
    if ext in DOCUMENT_EXTENSIONS:
        return True
    if ext in IMAGE_EXTENSIONS:
        return False
    # Unknown extension - try to read first bytes to detect
    try:
        with open(filepath, "rb") as f:
            header = f.read(12)
        # PNG: 89 50 4E 47
        if header[:4] == b"\x89PNG":
            return False
        # JPEG: FF D8 FF
        if header[:2] == b"\xff\xd8":
            return False
        # GIF: 47 49 46 38
        if header[:3] in (b"GIF",):
            return False
        # BMP: 42 4D
        if header[:2] == b"BM":
            return False
        # WEBP: 52 49 46 46 ... 57 45 42 50
        if header[:4] == b"RIFF" and header[8:12] == b"WEBP":
            return False
        return True
    except Exception:
        return True


async def convert_file(filename: str, folder: str = "") -> Optional[str]:
    """Convert a file in the inbox to markdown. Never raises - returns None on any error.
    'folder' determines where the note goes; for images it also sets img/ subfolder.
    """
    src = BRAIN_INBOX / filename
    if not src.exists():
        return None

    ext = src.suffix.lower()

    if ext in IMAGE_EXTENSIONS or _is_likely_image_bytes(src):
        md_folder = folder if folder else "img"
        img_subfolder = folder if folder else ""
    elif ext in {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".html", ".htm", ".txt", ".csv", ".json", ".xml", ".rtf", ".odt", ".epub"}:
        md_folder = folder if folder else "documentos"
        img_subfolder = ""
    else:
        md_folder = folder if folder else ""
        img_subfolder = ""

    md_filename = (md_folder + "/" if md_folder else "") + src.stem + ".md"
    dst = BRAIN_NOTES / md_filename
    dst.parent.mkdir(parents=True, exist_ok=True)

    try:
        if ext in IMAGE_EXTENSIONS or _is_likely_image_bytes(src):
            return await _convert_image(src, dst, filename, subfolder=img_subfolder)

        # --- AUDIO FILES: Skip ---
        if ext in {".mp3", ".wav", ".m4a", ".ogg", ".flac", ".wma", ".aac", ".webm"}:
            return None

        # --- DOCUMENT FILES: Native extraction per format (no MarkItDown to avoid image errors) ---
        text = None
        import subprocess, io, sys, logging
        if ext == ".pdf":
            try:
                proc = await asyncio.create_subprocess_exec(
                    "python", "-m", "pypdf", str(src),
                    stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
                )
                stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=30)
                text = stdout.decode("utf-8", errors="replace")
            except Exception:
                text = None
        if not text and ext == ".docx":
            try:
                from docx import Document
                doc = Document(str(src))
                text = "\n\n".join(p.text for p in doc.paragraphs)
            except Exception:
                text = None
        if not text and ext in {".pptx", ".ppt"}:
            try:
                from pptx import Presentation
                prs = Presentation(str(src))
                paras = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            paras.append(shape.text)
                text = "\n\n".join(paras)
            except Exception:
                text = None
        if not text and ext == ".html":
            try:
                from bs4 import BeautifulSoup
                with open(str(src), "r", encoding="utf-8", errors="replace") as _fh:
                    soup = BeautifulSoup(_fh.read(), "html.parser")
                text = soup.get_text(separator="\n").strip()
            except Exception:
                text = None
        if not text and ext in {".csv", ".txt", ".json", ".xml"}:
            try:
                text = src.read_text(encoding="utf-8", errors="replace")
            except Exception:
                text = None
        if not text:
            try:
                from markitdown import MarkItDown
                old_level = logging.root.level
                logging.root.setLevel(logging.CRITICAL)
                old_stderr = sys.stderr
                old_stdout = sys.stdout
                sys.stderr = io.StringIO()
                sys.stdout = io.StringIO()
                try:
                    md_converter = MarkItDown(llm=None, enable_builtins=False)
                    result = md_converter.convert(str(src))
                    text = result.text_content
                except Exception as md_e:
                    if "does not support image input" in str(md_e):
                        print(f"[Baul] MarkItDown image error (suppressed): {md_e}")
                    text = None
                finally:
                    sys.stderr = old_stderr
                    sys.stdout = old_stdout
                    logging.root.setLevel(old_level)
            except Exception:
                text = None

        if not text or text.strip() == "":
            if _is_likely_image_bytes(src):
                return await _convert_image(src, dst, filename)
            return None

        async with aiofiles.open(str(dst), "w", encoding="utf-8") as f:
            await f.write(text)

        src.unlink()
        return md_filename
    except Exception as e:
        print(f"[Baul] convert_file error for {filename}: {e}")
        return None


async def _convert_image(src: Path, dst: Path, original_filename: str, subfolder: str = "") -> Optional[str]:
    """Convert image file to markdown using OCR or a metadata fallback. Never raises."""
    loop = asyncio.get_event_loop()
    dst.parent.mkdir(parents=True, exist_ok=True)

    title = src.stem.replace("-", " ").replace("_", " ").title()

    try:
        meta = await loop.run_in_executor(None, _image_metadata_sync, str(src))
    except Exception:
        meta = None

    # Optimize image (resize + WebP) into img/[subfolder/]
    opt_result = await loop.run_in_executor(None, optimize_image_sync, str(src), 1920, 85, subfolder)
    if opt_result:
        opt_path, opt_meta = opt_result
    else:
        opt_path, opt_meta = None, None

    try:
        img_url = f"/api/images/{opt_path}" if opt_path else ""
        img_tag = f"\n![]({img_url})\n\n" if img_url else ""
        savings = f" | Ahorro: {opt_meta['savings_pct']}%" if opt_meta else ""

        if meta:
            dims = f"{opt_meta['width']}x{opt_meta['height']}" if opt_meta else f"{meta['width']}x{meta['height']}"
            fmt = f"webp" if opt_meta else meta['format']
            md = f"""# {title}

> Imagen procesada por Baul
> {datetime.now().strftime("%Y-%m-%d %H:%M")}
> Dimensiones: {dims} | Formato: {fmt}{savings}

{img_tag}---
"""
        else:
            md = f"""# {title}

> Imagen procesada por Baul
> {datetime.now().strftime("%Y-%m-%d %H:%M")}

{img_tag}---
"""
        try:
            ocr_text = await loop.run_in_executor(None, _ocr_image_sync, str(src))
            if ocr_text and ocr_text.strip():
                md += f"""## Texto extraído de la imagen

{ocr_text}

---

"""
            else:
                md += "*No se detectó texto en la imagen.*\n\n---\n\n"
        except Exception as ocr_e:
            md += f"*OCR: {ocr_e}*\n\n---\n\n"

        if opt_meta:
            md += f"""## Optimización

| Propiedad | Valor |
|-----------|-------|
| Original | {opt_meta['size'] / 1024:.1f} KB ({opt_meta.get('original_size', (opt_meta['width'], opt_meta['height']))}) |
| Optimizado | {opt_meta['optimized_size'] / 1024:.1f} KB ({opt_meta['width']}x{opt_meta['height']}) |
| Ahorro | {opt_meta['savings_pct']}% |
| Formato | WebP |
"""
        elif meta:
            md += f"""## Metadatos

| Propiedad | Valor |
|-----------|-------|
| Archivo original | {original_filename} |
| Ancho | {meta['width']} px |
| Alto | {meta['height']} px |
| Formato | {meta['format']} |
| Modo de color | {meta['mode']} |
| Fecha | {datetime.now().strftime('%Y-%m-%d %H:%M')} |
"""
        else:
            md += f"""**Archivo original:** {original_filename}
"""
        md += "\n---\n\n*Imagen procesada por Baul.*"
    except Exception as e:
        md = f"""# {title}

> Imagen procesada por Baul
> {datetime.now().strftime("%Y-%m-%d %H:%M")}

*Error al procesar: {e}*
"""

    try:
        async with aiofiles.open(str(dst), "w", encoding="utf-8") as f:
            await f.write(md)
    except Exception:
        return None

    try:
        src.unlink()
    except Exception:
        pass

    return dst.name


def is_image_file(filename: str) -> bool:
    """Check if a file is an image based on extension."""
    return Path(filename).suffix.lower() in IMAGE_EXTENSIONS


async def convert_all_inbox() -> list:
    """Convert all files in inbox to markdown."""
    results = []
    for f in sorted(BRAIN_INBOX.iterdir()):
        if f.is_file() and f.suffix.lower() not in (".md", ".markdown"):
            try:
                md_name = await convert_file(f.name)
                if md_name:
                    results.append(md_name)
            except Exception as e:
                results.append({"file": f.name, "error": str(e)})
    return results
